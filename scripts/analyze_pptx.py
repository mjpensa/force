import json
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

def analyze_pptx(file_path):
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return

    prs = Presentation(file_path)
    
    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    print(f"Slide Dimensions: {slide_width.inches} x {slide_height.inches} inches")
    
    templates = []
    
    for i, slide in enumerate(prs.slides):
        print(f"Analyzing Slide {i+1}...")
        elements = []
        
        for shape in slide.shapes:
            element = {
                "id": shape.shape_id,
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left.inches,
                "top": shape.top.inches,
                "width": shape.width.inches,
                "height": shape.height.inches,
                "rotation": shape.rotation
            }
            
            if shape.has_text_frame:
                text_frame = shape.text_frame
                element["text"] = text_frame.text
                
                # Analyze paragraphs for font info
                paragraphs = []
                for p in text_frame.paragraphs:
                    runs = []
                    for r in p.runs:
                        run_info = {
                            "text": r.text,
                            "bold": r.font.bold,
                            "italic": r.font.italic,
                            "size": r.font.size.pt if r.font.size else None,
                            "color": str(r.font.color.rgb) if r.font.color and hasattr(r.font.color, 'rgb') else None
                        }
                        runs.append(run_info)
                    paragraphs.append({"runs": runs, "alignment": str(p.alignment)})
                element["paragraphs"] = paragraphs
                
            elements.append(element)
            
        # Sort elements by top position to roughly order them
        elements.sort(key=lambda x: x['top'])
        
        templates.append({
            "slide_index": i + 1,
            "elements": elements
        })
        
    output = {
        "slide_width_inches": slide_width.inches,
        "slide_height_inches": slide_height.inches,
        "templates": templates
    }
    
    print(json.dumps(output, indent=2))

if __name__ == "__main__":
    analyze_pptx(r"c:\FORCE\PPT Templates_SHORT.pptx")
