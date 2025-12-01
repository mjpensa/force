from pptx import Presentation
import json

prs = Presentation(r'c:\FORCE\PPT Templates_SHORT3.pptx')
print(f'Slide Dimensions: {prs.slide_width.inches} x {prs.slide_height.inches} inches')
print(f'Total Slides: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    for shape in sorted(slide.shapes, key=lambda s: (s.top, s.left)):
        print(f'  [{shape.name}] type={shape.shape_type}')
        print(f'    pos: ({shape.left.inches:.2f}, {shape.top.inches:.2f}) size: {shape.width.inches:.2f} x {shape.height.inches:.2f}')
        if shape.has_text_frame:
            text = shape.text_frame.text[:60].replace('\n', ' ')
            print(f'    text: "{text}..."' if len(shape.text_frame.text) > 60 else f'    text: "{text}"')
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        color = str(r.font.color.rgb) if r.font.color and hasattr(r.font.color, 'rgb') and r.font.color.rgb else 'theme'
                        print(f'      font: {r.font.size.pt}pt, color={color}, bold={r.font.bold}')
                        break
    print()
